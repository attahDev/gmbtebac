import os
import logging
import hashlib
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from app.config import settings

logger = logging.getLogger(__name__)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

THEMES = {
    "dark": {
        "bg": RGBColor(0x0D, 0x0D, 0x1A),
        "bg_cover": RGBColor(0x08, 0x08, 0x18),
        "bg_slide": RGBColor(0x12, 0x12, 0x2B),
        "accent": RGBColor(0x4F, 0x8E, 0xFF),
        "accent2": RGBColor(0x00, 0xD4, 0xAA),
        "accent_text": RGBColor(0x00, 0xD4, 0xAA),
        "heading": RGBColor(0xFF, 0xFF, 0xFF),
        "subheading": RGBColor(0x00, 0xD4, 0xAA),
        "body": RGBColor(0xCC, 0xCC, 0xDD),
        "divider": RGBColor(0x2A, 0x2A, 0x4A),
        "label": RGBColor(0x4F, 0x8E, 0xFF),
        "slide_num": RGBColor(0xCC, 0xCC, 0xDD),
        "vector": RGBColor(0x4F, 0x8E, 0xFF),
        "vector2": RGBColor(0x00, 0xD4, 0xAA),
    },
    "light": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "bg_cover": RGBColor(0xF5, 0xF7, 0xFF),
        "bg_slide": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x2D, 0x6A, 0xFF),
        "accent2": RGBColor(0x00, 0xAA, 0x88),
        "accent_text": RGBColor(0x00, 0xAA, 0x88),
        "heading": RGBColor(0x11, 0x11, 0x33),
        "subheading": RGBColor(0x00, 0xAA, 0x88),
        "body": RGBColor(0x33, 0x33, 0x55),
        "divider": RGBColor(0xDD, 0xDD, 0xEE),
        "label": RGBColor(0x2D, 0x6A, 0xFF),
        "slide_num": RGBColor(0x99, 0x99, 0xAA),
        "vector": RGBColor(0x2D, 0x6A, 0xFF),
        "vector2": RGBColor(0x00, 0xAA, 0x88),
    },
    "corporate": {
        "bg": RGBColor(0x0A, 0x1A, 0x3A),
        "bg_cover": RGBColor(0x06, 0x10, 0x28),
        "bg_slide": RGBColor(0x0D, 0x1F, 0x44),
        "accent": RGBColor(0xC9, 0xA8, 0x4C),   # gold
        "accent2": RGBColor(0xE8, 0xC9, 0x7A),
        "accent_text": RGBColor(0xE8, 0xC9, 0x7A),
        "heading": RGBColor(0xFF, 0xFF, 0xFF),
        "subheading": RGBColor(0xC9, 0xA8, 0x4C),
        "body": RGBColor(0xCC, 0xD6, 0xEE),
        "divider": RGBColor(0x1E, 0x34, 0x60),
        "label": RGBColor(0xC9, 0xA8, 0x4C),
        "slide_num": RGBColor(0x88, 0x99, 0xBB),
        "vector": RGBColor(0xC9, 0xA8, 0x4C),
        "vector2": RGBColor(0xE8, 0xC9, 0x7A),
    },
    "minimal": {
        "bg": RGBColor(0xFA, 0xFA, 0xF8),
        "bg_cover": RGBColor(0xF0, 0xF0, 0xEC),
        "bg_slide": RGBColor(0xFA, 0xFA, 0xF8),
        "accent": RGBColor(0x22, 0x22, 0x22),
        "accent2": RGBColor(0x88, 0x88, 0x88),
        "accent_text": RGBColor(0x88, 0x88, 0x88),
        "heading": RGBColor(0x11, 0x11, 0x11),
        "subheading": RGBColor(0x55, 0x55, 0x55),
        "body": RGBColor(0x33, 0x33, 0x33),
        "divider": RGBColor(0xCC, 0xCC, 0xCC),
        "label": RGBColor(0x88, 0x88, 0x88),
        "slide_num": RGBColor(0xBB, 0xBB, 0xBB),
        "vector": RGBColor(0x22, 0x22, 0x22),
        "vector2": RGBColor(0xBB, 0xBB, 0xBB),
    },
    "bold": {
        "bg": RGBColor(0x00, 0x00, 0x00),
        "bg_cover": RGBColor(0x00, 0x00, 0x00),
        "bg_slide": RGBColor(0x0A, 0x0A, 0x0A),
        "accent": RGBColor(0xFF, 0x2D, 0x55),
        "accent2": RGBColor(0xFF, 0xCC, 0x00),
        "accent_text": RGBColor(0xFF, 0xCC, 0x00),
        "heading": RGBColor(0xFF, 0xFF, 0xFF),
        "subheading": RGBColor(0xFF, 0xCC, 0x00),
        "body": RGBColor(0xEE, 0xEE, 0xEE),
        "divider": RGBColor(0x22, 0x22, 0x22),
        "label": RGBColor(0xFF, 0x2D, 0x55),
        "slide_num": RGBColor(0x88, 0x88, 0x88),
        "vector": RGBColor(0xFF, 0x2D, 0x55),
        "vector2": RGBColor(0xFF, 0xCC, 0x00),
    },
    "gmbte": {
        "bg": RGBColor(0xFF, 0xFD, 0xF7),
        "bg_cover": RGBColor(0xFF, 0xF7, 0xE0),
        "bg_slide": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x00, 0x1F, 0x3F),
        "accent2": RGBColor(0xFF, 0xD7, 0x00),
        "accent_text": RGBColor(0xD7, 0x26, 0x3D),
        "heading": RGBColor(0x00, 0x1F, 0x3F),
        "subheading": RGBColor(0xD7, 0x26, 0x3D),
        "body": RGBColor(0x33, 0x41, 0x4F),
        "divider": RGBColor(0xE5, 0xE0, 0xCE),
        "label": RGBColor(0xD7, 0x26, 0x3D),
        "slide_num": RGBColor(0x9A, 0xA3, 0xAD),
        "vector": RGBColor(0xFF, 0xD7, 0x00),
        "vector2": RGBColor(0xD7, 0x26, 0x3D),
    },
}

THEME_PREVIEWS = [
    {"id": "dark",      "name": "Dark",      "description": "Dark navy with electric blue accents", "bg_color": "#0D0D1A", "accent_color": "#4F8EFF", "text_color": "#FFFFFF"},
    {"id": "light",     "name": "Light",     "description": "Clean white with blue accents",        "bg_color": "#FFFFFF", "accent_color": "#2D6AFF", "text_color": "#111133"},
    {"id": "corporate", "name": "Corporate", "description": "Deep navy with gold accents",           "bg_color": "#0A1A3A", "accent_color": "#C9A84C", "text_color": "#FFFFFF"},
    {"id": "minimal",   "name": "Minimal",   "description": "Off-white, clean typography",           "bg_color": "#FAFAF8", "accent_color": "#222222", "text_color": "#111111"},
    {"id": "bold",      "name": "Bold",      "description": "Black with high-contrast red & yellow", "bg_color": "#000000", "accent_color": "#FF2D55", "text_color": "#FFFFFF"},
    {"id": "gmbte",     "name": "GMBTE",     "description": "Cream with navy & yellow brand accents","bg_color": "#FFFDF7", "accent_color": "#001F3F", "text_color": "#001F3F"},
]

def _get_theme(theme_name: str) -> dict:
    return THEMES.get(theme_name, THEMES["gmbte"])


def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, left, top, width, height,
                 font_size, bold=False, color=None,
                 align=PP_ALIGN.LEFT, font_name="Calibri", wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = font_name
    return tb


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_oval(slide, left, top, size, color: RGBColor):
    shape = slide.shapes.add_shape(9, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_triangle(slide, left, top, width, height, color: RGBColor, rotation=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if rotation:
        shape.rotation = rotation
    return shape


def _add_dot_grid(slide, left, top, rows, cols, spacing, size, color: RGBColor):
    """A small grid of dots — a lightweight decorative 'vector' pattern
    distinct from the triangle+oval corner accent."""
    for r in range(rows):
        for c in range(cols):
            _add_oval(
                slide,
                left=left + Inches(c * spacing),
                top=top + Inches(r * spacing),
                size=Inches(size),
                color=color,
            )


def _add_chevron_stack(slide, left, top, color1: RGBColor, color2: RGBColor):
    """Stacked angled bars suggesting a chevron/arrow — a second decorative
    variant built from native shapes rather than an image."""
    _add_triangle(slide, left, top, Inches(0.5), Inches(1.4), color1, rotation=90)
    _add_triangle(slide, left + Inches(0.22), top, Inches(0.5), Inches(1.4), color2, rotation=90)


def _deck_variant_seed(deck_id: str) -> int:
    """Deterministic per-deck variant index, same pattern brand-identity-service
    uses for layout variation — same deck_id always renders the same way,
    different decks get visual variety without any randomness."""
    return int(hashlib.md5(deck_id.encode()).hexdigest(), 16)


def _add_vector_accents(slide, t: dict, variant="corner", deck_id: str = ""):
    v1 = t.get("vector")
    v2 = t.get("vector2")
    if not v1:
        return

    # Content/team/summary slides cycle through three decorative styles
    # based on the deck's own id, so a given deck is consistent slide-to-
    # slide feeling like the same design, while different decks don't all
    # look identical.
    style = _deck_variant_seed(deck_id) % 3 if deck_id else 0

    if variant == "corner":
        if style == 0:
            _add_triangle(slide, SLIDE_W - Inches(0.9), Inches(0),
                          Inches(0.9), Inches(0.9), v1, rotation=180)
            if v2:
                _add_oval(slide, Inches(0.15), SLIDE_H - Inches(0.35),
                          Inches(0.16), v2)
        elif style == 1:
            _add_dot_grid(slide, SLIDE_W - Inches(1.1), Inches(0.15),
                          rows=3, cols=4, spacing=0.22, size=0.08, color=v1)
        else:
            _add_chevron_stack(slide, SLIDE_W - Inches(0.9), Inches(0), v1, v2 or v1)
    elif variant == "cover":
        _add_triangle(slide, SLIDE_W - Inches(1.6), Inches(0),
                      Inches(1.6), Inches(1.6), v1, rotation=180)
        if v2:
            _add_oval(slide, SLIDE_W - Inches(2.0), SLIDE_H - Inches(2.0),
                      Inches(1.1), v2)
        if style == 1:
            _add_dot_grid(slide, Inches(0.4), SLIDE_H - Inches(1.0),
                          rows=3, cols=5, spacing=0.2, size=0.07, color=v2 or v1)


def _add_image(slide, image_source, left, top, width, height):
    try:
        if image_source.startswith("http"):
            r = requests.get(image_source, timeout=10)
            r.raise_for_status()
            img_stream = BytesIO(r.content)
            slide.shapes.add_picture(img_stream, left, top, width, height)
        else:
            if os.path.exists(image_source):
                slide.shapes.add_picture(image_source, left, top, width, height)
    except Exception as e:
        logger.warning(f"Failed to add image: {e}")


def _add_slide_number(slide, number: int, t: dict):
    _add_textbox(slide, str(number),
                 left=Inches(12.5), top=Inches(7.0),
                 width=Inches(0.6), height=Inches(0.4),
                 font_size=10, color=t["slide_num"], align=PP_ALIGN.RIGHT)


def _add_accent_bar(slide, t: dict):
    _add_rect(slide, Inches(0), Inches(0.12), SLIDE_W, Inches(0.06), t["accent"])



def _render_content_area(slide, heading, subheading, bullets, t, image_source=None):
    has_image = bool(image_source)
    text_width = Inches(7.5) if has_image else Inches(12.3)
    text_left = Inches(0.5)

    _add_textbox(slide, heading,
                 left=text_left, top=Inches(0.7),
                 width=text_width, height=Inches(1.1),
                 font_size=30, bold=True, color=t["heading"],
                 font_name="Calibri Light")
    if subheading:
        _add_textbox(slide, subheading,
                     left=text_left, top=Inches(1.85),
                     width=text_width, height=Inches(0.6),
                     font_size=15, color=t["subheading"])

    _add_rect(slide, text_left, Inches(2.55), text_width, Inches(0.04), t["divider"])

    for i, bullet in enumerate(bullets[:5]):
        _add_oval(slide,
                  left=text_left,
                  top=Inches(2.75) + Inches(i * 0.78) + Inches(0.1),
                  size=Inches(0.12), color=t["accent"])
        _add_textbox(slide, bullet,
                     left=text_left + Inches(0.28),
                     top=Inches(2.75) + Inches(i * 0.78),
                     width=text_width - Inches(0.3),
                     height=Inches(0.65),
                     font_size=14, color=t["body"])

    if has_image:
        _add_image(slide, image_source,
                   left=Inches(8.3), top=Inches(0.7),
                   width=Inches(4.5), height=Inches(6.3))



def _build_cover(prs, slide_data, company_name, tagline, t, image_source=None, deck_id=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg_cover"])

    if image_source:
        _add_image(slide, image_source,
                   left=Inches(0), top=Inches(0),
                   width=SLIDE_W, height=SLIDE_H)
        overlay = _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H,
                             RGBColor(0x00, 0x00, 0x00))
        overlay.fill.fore_color.theme_color = None
        try:
            overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
            overlay.fill.transparency = 0.45
        except Exception:
            pass

    _add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, t["accent"])
    _add_rect(slide, Inches(0.35), Inches(0), Inches(0.1), SLIDE_H, t["accent2"])
    _add_vector_accents(slide, t, variant="cover", deck_id=deck_id)

    _add_textbox(slide, company_name,
                 left=Inches(1.0), top=Inches(2.2),
                 width=Inches(9), height=Inches(1.5),
                 font_size=52, bold=True, color=t["heading"],
                 align=PP_ALIGN.LEFT, font_name="Calibri Light")

    _add_textbox(slide, tagline,
                 left=Inches(1.0), top=Inches(3.9),
                 width=Inches(9), height=Inches(0.8),
                 font_size=20, color=t["accent"])

    subheading = slide_data.get("subheading", "")
    if subheading:
        _add_textbox(slide, subheading,
                     left=Inches(1.0), top=Inches(4.8),
                     width=Inches(9), height=Inches(0.6),
                     font_size=14, color=t["body"])


def _build_content_slide(prs, slide_data, t, image_source=None, deck_id=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg_slide"])
    _add_accent_bar(slide, t)
    _add_vector_accents(slide, t, variant="corner", deck_id=deck_id)

    _add_textbox(slide, slide_data.get("title", "").upper(),
                 left=Inches(0.5), top=Inches(0.25),
                 width=Inches(4), height=Inches(0.35),
                 font_size=10, color=t["label"], bold=True)

    _render_content_area(
        slide,
        heading=slide_data.get("heading", ""),
        subheading=slide_data.get("subheading", ""),
        bullets=slide_data.get("bullets", []),
        t=t,
        image_source=image_source,
    )

    _add_slide_number(slide, slide_data.get("slide_number", 0), t)


def _build_team_slide(prs, slide_data, t, image_source=None, deck_id=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg_slide"])
    _add_accent_bar(slide, t)
    _add_vector_accents(slide, t, variant="corner", deck_id=deck_id)

    _add_textbox(slide, "TEAM",
                 left=Inches(0.5), top=Inches(0.25),
                 width=Inches(4), height=Inches(0.35),
                 font_size=10, color=t["label"], bold=True)

    has_image = bool(image_source)
    if has_image:
        _add_image(slide, image_source,
                   left=Inches(0.5), top=Inches(0.8),
                   width=Inches(4.5), height=Inches(6.0))
        text_left = Inches(5.3)
        text_width = Inches(7.5)
    else:
        text_left = Inches(0.5)
        text_width = Inches(12.3)

    _add_textbox(slide, slide_data.get("heading", ""),
                 left=text_left, top=Inches(0.7),
                 width=text_width, height=Inches(1.0),
                 font_size=28, bold=True, color=t["heading"],
                 font_name="Calibri Light")

    if slide_data.get("subheading"):
        _add_textbox(slide, slide_data["subheading"],
                     left=text_left, top=Inches(1.8),
                     width=text_width, height=Inches(0.55),
                     font_size=14, color=t["subheading"])

    _add_rect(slide, text_left, Inches(2.45), text_width, Inches(0.04), t["divider"])

    for i, bullet in enumerate(slide_data.get("bullets", [])[:4]):
        _add_oval(slide,
                  left=text_left,
                  top=Inches(2.65) + Inches(i * 0.95) + Inches(0.1),
                  size=Inches(0.12), color=t["accent2"])
        _add_textbox(slide, bullet,
                     left=text_left + Inches(0.28),
                     top=Inches(2.65) + Inches(i * 0.95),
                     width=text_width - Inches(0.3),
                     height=Inches(0.8),
                     font_size=14, color=t["body"])

    _add_slide_number(slide, slide_data.get("slide_number", 0), t)


def _build_summary_slide(prs, slide_data, t, deck_id=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg_slide"])
    _add_accent_bar(slide, t)
    _add_vector_accents(slide, t, variant="corner", deck_id=deck_id)

    _add_textbox(slide, "SUMMARY",
                 left=Inches(0.5), top=Inches(0.25),
                 width=Inches(4), height=Inches(0.35),
                 font_size=10, color=t["label"], bold=True)

    _add_textbox(slide, slide_data.get("heading", ""),
                 left=Inches(0.5), top=Inches(0.7),
                 width=Inches(12.3), height=Inches(1.0),
                 font_size=30, bold=True, color=t["heading"],
                 font_name="Calibri Light")

    if slide_data.get("subheading"):
        _add_textbox(slide, slide_data["subheading"],
                     left=Inches(0.5), top=Inches(1.8),
                     width=Inches(12.3), height=Inches(0.55),
                     font_size=15, color=t["subheading"])

    _add_rect(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.04), t["divider"])

    for i, bullet in enumerate(slide_data.get("bullets", [])[:5]):
        num_color = t["accent"] if i % 2 == 0 else t.get("accent_text", t["accent2"])
        _add_textbox(slide, str(i + 1),
                     left=Inches(0.5),
                     top=Inches(2.65) + Inches(i * 0.75),
                     width=Inches(0.35), height=Inches(0.6),
                     font_size=16, bold=True, color=num_color,
                     align=PP_ALIGN.CENTER)
        _add_textbox(slide, bullet,
                     left=Inches(1.0),
                     top=Inches(2.65) + Inches(i * 0.75),
                     width=Inches(11.8), height=Inches(0.65),
                     font_size=14, color=t["body"])

    _add_slide_number(slide, slide_data.get("slide_number", 0), t)


def _build_thankyou_slide(prs, slide_data, company_name, t, image_source=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg_cover"])

    if image_source:
        _add_image(slide, image_source,
                   left=Inches(0), top=Inches(0),
                   width=SLIDE_W, height=SLIDE_H)

    _add_rect(slide, Inches(0), Inches(6.9), SLIDE_W, Inches(0.6), t["accent"])

    _add_textbox(slide, slide_data.get("heading", "Thank You"),
                 left=Inches(1.5), top=Inches(2.0),
                 width=Inches(10), height=Inches(1.8),
                 font_size=54, bold=True, color=t["heading"],
                 align=PP_ALIGN.CENTER, font_name="Calibri Light")

    _add_textbox(slide, company_name,
                 left=Inches(1.5), top=Inches(3.9),
                 width=Inches(10), height=Inches(0.7),
                 font_size=22, color=t["accent"],
                 align=PP_ALIGN.CENTER)

    subheading = slide_data.get("subheading", "")
    if subheading:
        _add_textbox(slide, subheading,
                     left=Inches(1.5), top=Inches(4.7),
                     width=Inches(10), height=Inches(0.6),
                     font_size=15, color=t["body"],
                     align=PP_ALIGN.CENTER)


def build_pptx(slides_data: dict, deck_id: str,
               theme_name: str = "gmbte", slide_images: dict = None) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    t = _get_theme(theme_name)
    images = slide_images or {}

    company_name = slides_data.get("company_name", "Company")
    tagline = slides_data.get("tagline", "")
    slides = slides_data.get("slides", [])

    for slide_data in slides:
        num = slide_data.get("slide_number")
        title = slide_data.get("title", "")
        image = images.get(num)

        if title == "Cover":
            _build_cover(prs, slide_data, company_name, tagline, t, image, deck_id=deck_id)
        elif title == "Team":
            _build_team_slide(prs, slide_data, t, image, deck_id=deck_id)
        elif title == "Summary":
            _build_summary_slide(prs, slide_data, t, deck_id=deck_id)
        elif title == "Thank You":
            _build_thankyou_slide(prs, slide_data, company_name, t, image)
        else:
            _build_content_slide(prs, slide_data, t, image, deck_id=deck_id)

    os.makedirs(settings.MEDIA_DIR, exist_ok=True)
    file_path = os.path.join(settings.MEDIA_DIR, f"{deck_id}.pptx")
    prs.save(file_path)
    logger.info(f"PPTX saved: {file_path} (theme={theme_name}, slides={len(slides)})")
    return file_path
